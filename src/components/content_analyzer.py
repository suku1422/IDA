import streamlit as st

from src.openai_client import get_openai_response

def analyze_content():
    st.header("Step 2: Analyze Raw Content")
    context_summary = st.session_state.get("context_summary", "")
    uploaded_files = st.file_uploader(
        "Upload raw content files (PDF, DOCX, TXT, PPTX, XLSX)", 
        type=["pdf", "docx", "txt", "pptx", "xlsx"], 
        accept_multiple_files=True
    )

    if uploaded_files:
        raw_text = ""
        try:
            for uploaded_file in uploaded_files:
                if uploaded_file.type == "application/pdf":
                    from PyPDF2 import PdfReader
                    reader = PdfReader(uploaded_file)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            raw_text += page_text + "\n"
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    import docx
                    doc = docx.Document(uploaded_file)
                    for para in doc.paragraphs:
                        raw_text += para.text + "\n"
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    from pptx import Presentation
                    prs = Presentation(uploaded_file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                raw_text += shape.text + "\n"
                elif uploaded_file.type in [
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "application/vnd.ms-excel"
                ]:
                    import pandas as pd
                    xls = pd.ExcelFile(uploaded_file)
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(uploaded_file, sheet_name=sheet_name)
                        raw_text += f"\n--- Sheet: {sheet_name} ---\n"
                        raw_text += df.to_string(index=False) + "\n"
                elif uploaded_file.type == "text/plain":
                    raw_text += uploaded_file.getvalue().decode("utf-8") + "\n"
        except Exception as e:
            st.error(f"Error reading uploaded file: {e}")
            return

        if "analysis_done" not in st.session_state or st.session_state.raw_text != raw_text:
            st.session_state.raw_text = raw_text
            context_summary = st.session_state.get("context_summary", "No context summary available.")
            prompt = (
                f"Analyze the following instructional design context and raw content to identify content gaps.\n\n"
                f"Here is the instructional design context:\n{context_summary}\n\n"
                f"Here is the raw content:\n{raw_text}\n\n"
                f"Identify any content gaps in the raw content based on the provided context. Don't make it very elaborate and focus on the duration indicated by the user in {context_summary}. "
                f"List the missing topics or areas that need to be covered in the course."
            )
            with st.spinner("Analyzing content gaps..."):
                analysis = get_openai_response(prompt)
                st.session_state.analysis = analysis
                st.session_state.analysis_done = True
                st.session_state.uploaded_content = raw_text

        st.subheader("Content Gap Analysis")
        st.write(st.session_state.analysis)

        decision = st.radio(
            "How would you like to address the content gaps?",
            options=("Generate content to fill gaps", "Provide additional sources", "No action needed"),
            index=2
        )

        if decision == "Generate content to fill gaps":
            filled_prompt = (
                f"Based on the identified content gaps below, generate the necessary content to fill these gaps.\n\n"
                f"**Content Gaps:**\n{st.session_state.analysis}\n\n"
                f"Provide the additional content required to cover these areas effectively."
            )
            with st.spinner("Generating additional content..."):
                filled_content = get_openai_response(filled_prompt)
                if filled_content:
                    st.session_state.filled_content = filled_content
                    st.session_state.generated_additional_content = filled_content
                    st.success("Content gaps have been filled with generated material.")
                    if st.button("Continue to Step 3"):
                        st.session_state.step = 3
                        st.rerun()

        elif decision == "Provide additional sources":
            more_files = st.file_uploader(
                "Upload additional files to improve coverage",
                type=["pdf", "docx", "txt"],
                accept_multiple_files=True,
                key="more_sources"
            )
            if more_files:
                uploaded_files.extend(more_files)
                st.warning("Please click the 'Analyze Again' button below to include the new files.")
                if st.button("Analyze Again"):
                    st.rerun()

        elif decision == "No action needed":
            if st.button("Continue to Step 3"):
                st.session_state.step = 3
                st.rerun()
    else:
        st.info("Please upload at least one raw content file to begin analysis.")

def provide_feedback(analysis):
    feedback_message = "Analysis Results:\n"
    if analysis["content_gaps"]:
        feedback_message += "Identified Content Gaps:\n"
        for gap in analysis["content_gaps"]:
            feedback_message += f"- {gap}\n"
    else:
        feedback_message += "No content gaps identified.\n"

    feedback_message += analysis["feedback"]
    return feedback_message